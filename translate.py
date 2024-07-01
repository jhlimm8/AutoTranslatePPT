from pptx import Presentation
import re
import argparse

def is_chinese(text):
    """Check if the text contains Chinese characters."""
    return bool(re.search('[\u4e00-\u9fff]', text))

def translate_chinese_text(text, api_key, api_client, api_model):
    """Translate Chinese text using OpenAI's GPT-4 model."""
    if api_client == "OpenAI":
        from openai import OpenAI
        client = OpenAI(
            api_key= api_key
        )
    elif api_client == "TogetherAI":
        from together import Together
        client = Together(
            api_key= api_key
        )
    elif api_client == "Groq":
        from groq import Groq
        client = Groq(
            api_key= api_key
        )
    try:
        response = client.chat.completions.create(
            model=api_model,
            messages=[
                {"role": "system", "content": "You are a translator. Translate the following Chinese text to English. Only return the translated text."},
                {"role": "user", "content": text}
            ],
            max_tokens=1024
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"Error in translation: {e}")
        return text  # Return original text if translation fails

def process_slide(slide, api_key, api_client, api_model):
    """Process a single slide, translating Chinese text in text boxes."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if is_chinese(paragraph.text):
                    # Translate the text
                    translated_text = translate_chinese_text(paragraph.text, api_key, api_client, api_model)
                    print(f"Original text: {paragraph.text}")
                    print(f"Translated text: {translated_text}")
                    # Replace the original text with the translated text
                    paragraph.text = translated_text

def process_presentation(input_path, output_path, api_key, api_client, api_model):
    """Process the entire presentation and save the modified version."""
    prs = Presentation(input_path)
    
    for slide in prs.slides:
        process_slide(slide, api_key, api_client, api_model)
    
    prs.save(output_path)
    print(f"Modified presentation saved to {output_path}")

def main():
    parser = argparse.ArgumentParser(description="Translate Chinese text in PowerPoint slides to English.")
    parser.add_argument("api_key", help="API key")
    parser.add_argument("input_file", help="Path to the input PowerPoint file")
    parser.add_argument("output_file", help="Path to save the translated PowerPoint file")
    parser.add_argument("api_client", help="API client e.g. OpenAI, TogetherAI, Groq")
    parser.add_argument("api_model", help="API model e.g. GPT-4o, Qwen/Qwen2-72B-Instruct, llama3-70b-8192")
    
    args = parser.parse_args()
    process_presentation(args.input_file, args.output_file, args.api_key, args.api_client, args.api_model)

if __name__ == "__main__":
    main()